from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from analysis.cluster_explanation import ClusterExplanation
from analysis.clustering import ClusteringResult
from analysis.embedding_comparison import EmbeddingComparisonResult


def generate_ml_pptx_report(
    df: pd.DataFrame,
    embedding_result: EmbeddingComparisonResult,
    clustering_result: ClusteringResult,
    explanations: list[ClusterExplanation],
    image_paths: dict[str, str],
    output_path: str,
) -> str:
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    _title_slide(prs, "ML-анализ структуры данных", datetime.now().strftime("%Y-%m-%d %H:%M"))
    _bullets_slide(
        prs,
        "Общая информация",
        [
            f"Строк: {len(df)}",
            f"Колонок: {df.shape[1]}",
            f"Использованные признаки: {', '.join(embedding_result.used_columns[:12])}",
            f"Sample: {'да' if embedding_result.sampled or clustering_result.sampled else 'нет'}",
        ],
    )
    _image_slide(prs, "PCA / UMAP / t-SNE", image_paths.get("embeddings"), "Сравнение методов снижения размерности.")
    _image_slide(
        prs,
        "Кластеризация",
        image_paths.get("clusters"),
        f"Метод: {clustering_result.method}; кластеры: {clustering_result.n_clusters}; silhouette: {clustering_result.silhouette_score}",
    )
    _cluster_table_slide(prs, explanations)
    _bullets_slide(prs, "Рекомендации", [rec for item in explanations for rec in item.recommendations][:10])
    _bullets_slide(
        prs,
        "Ограничения анализа",
        [
            "UMAP/t-SNE являются методами визуализации.",
            "Кластеры требуют бизнес-проверки.",
            "Результат зависит от качества данных и набора числовых признаков.",
        ],
    )
    prs.save(path)
    return str(path)


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.clear()
    for bullet in bullets or ["Нет рекомендаций."]:
        p = frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)


def _image_slide(prs: Presentation, title: str, image_path: str | None, note: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(0.7), Inches(1.2), width=Inches(8.6))
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(8.6), Inches(0.5))
    box.text_frame.text = note


def _cluster_table_slide(prs: Presentation, explanations: list[ClusterExplanation]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Объяснение кластеров"
    rows = max(2, len(explanations) + 1)
    table = slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.2), Inches(9), Inches(4.8)).table
    for index, header in enumerate(["cluster_id", "size", "share", "top_features"]):
        table.cell(0, index).text = header
    for row, item in enumerate(explanations, start=1):
        table.cell(row, 0).text = str(item.cluster_id)
        table.cell(row, 1).text = str(item.size)
        table.cell(row, 2).text = f"{item.share_percent:.1f}%"
        table.cell(row, 3).text = ", ".join(item.top_features.keys())
